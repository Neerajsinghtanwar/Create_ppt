from pptx import Presentation
from pptx.util import Cm, Inches, Pt
from wand.image import Image
import os

# create presentation
prs = Presentation()

# choosing slide layout
sl = prs.slide_layouts[8]

num_of_slides = 5
slides = {}
images = {}
logo = 'images/nike_black.png'

# create directory for composite images
try:
    path = os.path.join(os.getcwd(), "logo_images")
    os.mkdir(path)

except FileExistsError:
    print("File is already exists")

for n in range(1,num_of_slides+1):
    # create slides
    slides["slide%s" %n] = prs.slides.add_slide(sl)
    # create image path
    images["image%s" % n] = 'images/image%s.jpg' % n

x = 0
for i, j in zip(images.values(), slides.values()):
    x = x+1
    # set logo in image
    print("set logo in Image%s" % x)
    with Image(filename=i) as img:
        img.resize(500, 400)
        with Image(filename=logo) as logo_img:
            logo_img.resize(170, 60)
            img.composite(logo_img, left=5, top=5)
            loc = "logo_images/logo_image%s.jpg" % x
        img.save(filename=loc)
        # adding image in slides
        j.shapes.add_picture(loc, Inches(1.5), Inches(1.5))
        # adding title in slides
        left = Inches(1.4)
        top = Inches(0.5)
        width = height = Inches(7)
        p = j.shapes.add_textbox(left, top, width, height).text_frame.add_paragraph()
        p.text = "This ppt file is created by Python"
        p.font.size = Pt(30)

# saving file
prs.save("my_ppt.pptx")